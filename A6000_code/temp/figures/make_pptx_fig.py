# -*- coding: utf-8 -*-
"""Professional framework figure as editable PPTX (python-pptx):
gradient rounded boxes, exact formulas, colour-coded EDM + SDA paths."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

# palette (RGB)
FC = dict(DATA=RGBColor(0xDC,0xE9,0xF7), CONV=RGBColor(0xE3,0xF0,0xE6), DOWN=RGBColor(0xCD,0xE3,0xF5),
          UP=RGBColor(0xC8,0xEA,0xD3), NOISE=RGBColor(0xF5,0xE6,0xD3), DA=RGBColor(0xFB,0xE3,0xDE),
          LOSS=RGBColor(0xED,0xE3,0xF6))
EC = dict(DATA=RGBColor(0x3E,0x6B,0x9E), CONV=RGBColor(0x2F,0x7D,0x46), DOWN=RGBColor(0x1A,0x5A,0x9E),
          UP=RGBColor(0x1E,0x7A,0x7A), NOISE=RGBColor(0xB0,0x6A,0x1B), DA=RGBColor(0xC0,0x39,0x2B),
          LOSS=RGBColor(0x6A,0x3D,0x9A), GREY=RGBColor(0x44,0x44,0x44))

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

def add_box(slide, x, y, w, h, kind, text="", fs=10, bold=False, tc=None, fill=True, line_w=1.2):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.adjustments[0] = 0.08
    if fill:
        shp.fill.solid(); shp.fill.fore_color.rgb = FC[kind]
    else:
        shp.fill.background()
    shp.line.color.rgb = EC[kind]; shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    tf = shp.text_frame; tf.word_wrap = True
    tf.margin_left = Pt(2); tf.margin_right = Pt(2); tf.margin_top = Pt(1); tf.margin_bottom = Pt(1)
    if text:
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = text
        r.font.size = Pt(fs); r.font.bold = bold
        r.font.color.rgb = tc if tc else EC[kind]
        r.font.name = "Helvetica"
    return shp

def add_arrow(slide, x1, y1, x2, y2, color=RGBColor(0x44,0x44,0x44), lw=1.2):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = color; conn.line.width = Pt(lw)
    # arrowhead
    ln = conn.line._get_or_add_ln()
    head = ln.makeelement(qn('a:tailEnd'), {'type':'arrow', 'w':'med', 'len':'med'})
    ln.append(head)
    return conn

# ===== title =====
add_box(slide, 0.4, 0.2, 12.5, 0.55, "DATA", text="SDA-Diff: generative data assimilation for dense storm-surge forecasting",
        fs=15, bold=True, tc=EC["DATA"])

# ===== (a) INPUT (left) =====
add_box(slide, 0.4, 1.0, 2.9, 5.9, "DATA", text="", fill=True)
add_box(slide, 0.55, 1.15, 2.6, 0.4, "DATA", text="Input", fs=11, bold=True, tc=EC["DATA"])
add_box(slide, 0.55, 1.7, 2.6, 0.75, "DATA", text="sparse in-situ gauges\n+ valid mask", fs=8)
add_box(slide, 0.55, 2.55, 2.6, 0.75, "DATA", text="ERA5  msl·u10·v10", fs=8)
add_box(slide, 0.55, 3.4, 2.6, 0.75, "DATA", text="GTSM surge", fs=8)
add_box(slide, 0.55, 4.3, 2.6, 0.75, "DATA", text="c = [B,T,6,H,W]  +  lead L", fs=8, bold=True)
add_arrow(slide, 0.55, 5.2, 3.35, 5.6, color=EC["GREY"])

# ===== MODEL (centre) =====
add_box(slide, 3.5, 1.0, 6.3, 5.9, "CONV", text="", fill=True, line_w=1.6)
add_box(slide, 3.65, 1.15, 6.0, 0.4, "CONV", text="Model — EDM generation + SDA assimilation", fs=11, bold=True, tc=EC["CONV"])

# training path (top)
add_box(slide, 3.7, 1.75, 5.95, 2.1, "CONV", text="", fill=False, line_w=1.0)
add_box(slide, 3.85, 1.85, 5.7, 0.35, "CONV", text="Training — learn to recover the surge field from noise", fs=9, bold=True, tc=EC["CONV"], fill=False)
add_box(slide, 3.9, 2.3, 1.2, 0.7, "DATA", text="x₀ target", fs=8)
add_box(slide, 5.2, 2.3, 1.7, 0.7, "NOISE", text="xₜ = x₀ + σε\n(noise / drop)", fs=7.5)
add_arrow(slide, 5.1, 2.65, 5.2, 2.65, color=EC["NOISE"])
add_box(slide, 7.0, 2.3, 2.4, 0.7, "CONV", text="D_θ(xₜ,σ,c,L)  U-Net\nconv·down·up·attn", fs=7.5, bold=True)
add_arrow(slide, 6.9, 2.65, 7.0, 2.65)

# inference path (bottom)
add_box(slide, 3.7, 4.1, 5.95, 2.65, "DA", text="", fill=False, line_w=1.4)
add_box(slide, 3.85, 4.2, 5.7, 0.35, "DA", text="Inference — correct the sampling path with new observations", fs=9, bold=True, tc=EC["DA"], fill=False)
add_box(slide, 3.9, 4.65, 1.1, 0.7, "DATA", text="noise xₜ", fs=8)
add_box(slide, 5.1, 4.65, 1.3, 0.7, "CONV", text="denoise\n(EDM Heun)", fs=7.5)
add_arrow(slide, 5.0, 5.0, 5.1, 5.0)
add_box(slide, 6.5, 4.65, 1.3, 0.7, "CONV", text="Tweedie x̂₀", fs=7.5)
add_arrow(slide, 6.4, 5.0, 6.5, 5.0)
add_box(slide, 7.9, 4.65, 1.7, 0.7, "DATA", text="obs. y, A, R\n(gauge pixels)", fs=7.5)
add_box(slide, 5.0, 5.55, 4.5, 1.0, "DA", text="SDA:  ∇log p(xₜ|y,c) = s_θ + ∇log N(y|Ax̂₀,R)", fs=8.5, bold=True, tc=EC["DA"], line_w=1.8)
add_arrow(slide, 7.1, 4.65, 7.25, 5.55, color=EC["DA"])
add_arrow(slide, 8.7, 5.0, 7.25, 5.55, color=EC["DA"])

# ===== (c) OUTPUT (right) =====
add_box(slide, 10.0, 1.0, 3.0, 5.9, "DATA", text="", fill=True)
add_box(slide, 10.15, 1.15, 2.7, 0.4, "DATA", text="Output (ensemble)", fs=11, bold=True, tc=EC["DATA"])
add_box(slide, 10.2, 2.2, 2.6, 0.7, "CONV", text="posterior {x₀⁽ⁱ⁾} (N samples)", fs=8)
add_box(slide, 10.2, 3.3, 2.6, 1.0, "UP", text="mean · quantiles\nP(surge > h)", fs=8.5)
add_box(slide, 10.2, 4.7, 2.6, 0.8, "UP", text="dense surge forecast\n(uncertainty-aware)", fs=8, bold=True)
add_arrow(slide, 9.8, 2.6, 10.0, 2.6)
add_arrow(slide, 7.25, 5.9, 10.0, 3.6, color=EC["DA"])

# ===== loss (bottom) =====
add_box(slide, 3.7, 7.0, 6.0, 0.45, "LOSS", text="L = E[λ(σ) ‖ D_θ(x₀+σε, σ, c, L) − x₀ ‖²]   (masked, weighted)", fs=9, bold=True, tc=EC["LOSS"])
add_arrow(slide, 6.5, 3.85, 6.5, 7.0, color=EC["LOSS"])

prs.save("temp/figures/fig1_framework.pptx")
print("saved fig1_framework.pptx (editable)")
