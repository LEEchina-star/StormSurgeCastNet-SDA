# -*- coding: utf-8 -*-
"""Gulf case PPTX (2 slides): overview figures + 12-step GTSM diffusion series."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
NAVY = RGBColor(0x1A, 0x5A, 0x9E); GREY = RGBColor(0x44, 0x44, 0x44)

def title(s, text, x, y, w, h, fs=13, color=NAVY):
    box = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = text; r.font.size = Pt(fs); r.font.bold = True; r.font.color.rgb = color
    return box

def caption(s, lines, x, y, w, h, fs=8.5):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run(); r.text = line; r.font.size = Pt(fs); r.font.color.rgb = GREY

def pic(s, path, x, y, w, h):
    return s.shapes.add_picture(path, Inches(x), Inches(y), Inches(w), Inches(h))

# ================= SLIDE 1: overview =================
s = prs.slides.add_slide(prs.slide_layouts[6])
title(s, "SDA-Diff case study \u2014 Gulf of Mexico (st528, 28.99\u00b0N / 89.14\u00b0W; 5 GESLA-3 gauges in ROI, max obs 1.40 m)",
      0.3, 0.12, 12.7, 0.5, fs=14)
title(s, "(a) Geographic setting", 0.35, 0.8, 4.4, 0.35, fs=12)
pic(s, "temp/figures/gulf3_map.png", 0.35, 1.2, 4.5, 3.9)
caption(s, ["ROI-internal GESLA-3 stations (blue dots); selected station",
            "st528 (red star). Dashed box = 256\u00d7256 px ROI at 0.025\u00b0/px."],
        0.35, 5.15, 4.5, 1.6)
title(s, "(b) Input  c = [B, T=12, 6, H, W]", 5.05, 0.8, 4.3, 0.35, fs=12)
pic(s, "temp/figures/gulf3_inputs.png", 5.05, 1.2, 3.9, 3.9)
caption(s, ["12 hourly frames: sparse GESLA-3 surge (m), ERA5 MSLP (hPa),",
            "ERA5 wind speed (m/s), GTSM surge (m), rasterised onto ROI."],
        5.05, 5.15, 4.0, 1.7)
title(s, "(c) Generative denoising (EDM)", 9.15, 0.8, 3.9, 0.35, fs=12)
pic(s, "temp/figures/gulf3_denoise.png", 9.15, 1.2, 3.85, 1.5)
title(s, "(d) SDA data assimilation", 9.15, 2.85, 3.9, 0.35, fs=12)
pic(s, "temp/figures/gulf_assim.png", 9.15, 3.25, 3.85, 1.5)
caption(s, ["(c) clean \u2192 noisy \u2192 denoised surge field (m).",
            "(d) prior \u2192 observation (dense GTSM surge distribution) \u2192 posterior:",
            "SDA likelihood guidance assimilates GESLA-3 into a dense field."],
        9.15, 4.85, 3.9, 2.3)

# ================= SLIDE 2: 12-step GTSM series + assimilation =================
s2 = prs.slides.add_slide(prs.slide_layouts[6])
title(s2, "GTSM surge (land-masked) \u2014 12 hourly frames (t-0h \u2026 t-11h) + SDA assimilation",
      0.3, 0.15, 12.7, 0.5, fs=14)
for j, (f, lbl) in enumerate([("gulf_gtsm_input.png", "(a) input 12 frames (clean)"),
                              ("gulf_gtsm_noisy.png", "(b) + Gaussian noise 12 frames"),
                              ("gulf_assim.png", "(c) assimilation: prior \u2192 obs \u2192 posterior")]):
    x = 0.3 + j * 4.32
    title(s2, lbl, x, 0.8, 4.2, 0.35, fs=11)
    pic(s2, f"temp/figures/{f}", x, 1.15, 4.15, 3.2)
caption(s2, ["(a,b) GTSM surge (coarse, 72 stations) with land masked out (valid mask), north-up; 12 frames t-0h (latest) \u2192 t-11h (earliest).",
            "(c) SDA assimilation: prior p(x|c) \u2192 observation (GESLA-3 gauges inside ROI) \u2192 posterior p(x|c,y).",
            "All colourbars in metres (m)."],
        0.3, 4.6, 12.7, 0.9, fs=9)

prs.save("temp/figures/gulf_case.pptx")
print("saved gulf_case.pptx (2 slides)")
