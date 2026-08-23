# -*- coding: utf-8 -*-
"""Compose Gulf case (v3 figures) into a single editable PPTX."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
s = prs.slides.add_slide(prs.slide_layouts[6])

NAVY = RGBColor(0x1A, 0x5A, 0x9E); GREY = RGBColor(0x44, 0x44, 0x44)

def title(text, x, y, w, h, fs=13, color=NAVY):
    box = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = text; r.font.size = Pt(fs); r.font.bold = True; r.font.color.rgb = color
    return box

def caption(lines, x, y, w, h, fs=8.5):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run(); r.text = line; r.font.size = Pt(fs); r.font.color.rgb = GREY

def pic(path, x, y, w, h):
    return s.shapes.add_picture(path, Inches(x), Inches(y), Inches(w), Inches(h))

# main title
title("SDA-Diff case study \u2014 Gulf of Mexico (st528, 28.99\u00b0N / 89.14\u00b0W; 5 GESLA-3 gauges in ROI, max obs 1.40 m)",
      0.3, 0.12, 12.7, 0.5, fs=14)

# ---- left: geographic map ----
title("(a) Geographic setting", 0.35, 0.8, 4.4, 0.35, fs=12)
pic("temp/figures/gulf3_map.png", 0.35, 1.2, 4.5, 3.9)
caption(["ROI-internal GESLA-3 stations (blue dots); selected",
         "station (red star). Dashed box = 256\u00d7256 px ROI at",
         "0.025\u00b0/px (\u2248 6.4\u00b0 \u2248 640 km)."],
        0.35, 5.15, 4.5, 1.6)

# ---- middle: input subplots ----
title("(b) Input  c = [B, T=12, 6, H, W]", 5.05, 0.8, 4.3, 0.35, fs=12)
pic("temp/figures/gulf3_inputs.png", 5.05, 1.2, 3.9, 3.9)
caption(["12 hourly frames of: sparse GESLA-3 in-situ surge (m),",
         "ERA5 mean sea-level pressure (hPa), ERA5 wind speed",
         "(m/s), and GTSM surge (m), rasterised onto the ROI grid."],
        5.05, 5.15, 4.0, 1.7)

# ---- right: process diagrams ----
title("(c) Generative denoising (EDM)", 9.15, 0.8, 3.9, 0.35, fs=12)
pic("temp/figures/gulf3_denoise.png", 9.15, 1.2, 3.85, 1.5)
title("(d) SDA data assimilation", 9.15, 2.85, 3.9, 0.35, fs=12)
pic("temp/figures/gulf3_assim.png", 9.15, 3.25, 3.85, 1.5)
caption(["(c) clean \u2192 noisy \u2192 denoised surge field (m).",
         "(d) prior p(x|c) \u2192 sparse GESLA-3 observation y (5 gauges, red star = st528, 0.53 m) \u2192",
         "posterior p(x|c,y): when gauge observations exist, SDA likelihood",
         "guidance assimilates them into a dense, high-resolution field."],
        9.15, 4.85, 3.9, 2.3)

prs.save("temp/figures/gulf_case.pptx")
print("saved gulf_case.pptx (v3 figures)")
