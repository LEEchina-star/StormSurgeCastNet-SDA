# -*- coding: utf-8 -*-
"""Compose Gulf case into a single PPTX: map + inputs + process diagrams."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
s = prs.slides.add_slide(prs.slide_layouts[6])

def title(text, x, y, w, h, fs=14, color=RGBColor(0x1A,0x5A,0x9E)):
    box = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = text; r.font.size = Pt(fs); r.font.bold = True; r.font.color.rgb = color
    return box

def pic(path, x, y, w, h):
    return s.shapes.add_picture(path, Inches(x), Inches(y), Inches(w), Inches(h))

# main title
title("SDA-Diff: multi-modal input and generative assimilation — Gulf of Mexico case (id=528, surge 8.25)",
      0.3, 0.15, 12.7, 0.5, fs=15)

# ---- left: geographic map ----
title("(a) Geographic setting", 0.35, 0.85, 4.4, 0.4, fs=12)
pic("temp/figures/gulf_map.png", 0.35, 1.3, 4.55, 3.2)
# caption
tb = s.shapes.add_textbox(Inches(0.35), Inches(4.55), Inches(4.5), Inches(2.2))
tf = tb.text_frame; tf.word_wrap = True
for i, line in enumerate(["Gulf of Mexico storm-surge gauges (blue dots);",
                          "selected gauge (red star) with its 256×256 ROI",
                          "(dashed box, 0.025°/px ≈ 640 km)."]):
    p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
    r = p.add_run(); r.text = line; r.font.size = Pt(9); r.font.color.rgb = RGBColor(0x44,0x44,0x44)

# ---- middle: input subplots ----
title("(b) Input channels c=[B,T,6,H,W]", 5.1, 0.85, 4.3, 0.4, fs=12)
pic("temp/figures/gulf_inputs.png", 5.1, 1.3, 4.0, 4.0)
tb = s.shapes.add_textbox(Inches(5.1), Inches(5.35), Inches(4.2), Inches(1.5))
tf = tb.text_frame; tf.word_wrap = True
for i, line in enumerate(["Sparse in-situ gauges, validity mask, ERA5",
                          "pressure & 10-m wind, and coarse GTSM surge, all",
                          "rasterised onto the ROI grid (12 hourly frames)."]):
    p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
    r = p.add_run(); r.text = line; r.font.size = Pt(9); r.font.color.rgb = RGBColor(0x44,0x44,0x44)

# ---- right: process diagrams ----
title("(c) Generative denoising (EDM)", 9.5, 0.85, 3.5, 0.4, fs=12)
pic("temp/figures/gulf_denoise.png", 9.5, 1.3, 3.5, 1.28)
title("(d) SDA data assimilation", 9.5, 2.75, 3.5, 0.4, fs=12)
pic("temp/figures/gulf_assim.png", 9.5, 3.2, 3.5, 1.28)
tb = s.shapes.add_textbox(Inches(9.5), Inches(4.55), Inches(3.5), Inches(2.4))
tf = tb.text_frame; tf.word_wrap = True
for i, line in enumerate(["(c) clean → noisy → denoised surge field;",
                          "(d) prior (no obs) → observation (red star) →",
                          "posterior after likelihood-guided assimilation.",
                          "The generative model densifies the sparse gauges",
                          "into a high-resolution, uncertainty-aware field."]):
    p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
    r = p.add_run(); r.text = line; r.font.size = Pt(9); r.font.color.rgb = RGBColor(0x44,0x44,0x44)

prs.save("temp/figures/gulf_case.pptx")
print("saved gulf_case.pptx")
