# -*- coding: utf-8 -*-
"""Compose: SD background image + precise framework (PPTX), 3 variants."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from PIL import Image

FC = dict(DATA=RGBColor(0xDC,0xE9,0xF7), CONV=RGBColor(0xE3,0xF0,0xE6), DOWN=RGBColor(0xCD,0xE3,0xF5),
          UP=RGBColor(0xC8,0xEA,0xD3), NOISE=RGBColor(0xF5,0xE6,0xD3), DA=RGBColor(0xFB,0xE3,0xDE),
          LOSS=RGBColor(0xED,0xE3,0xF6))
EC = dict(DATA=RGBColor(0x3E,0x6B,0x9E), CONV=RGBColor(0x2F,0x7D,0x46), DOWN=RGBColor(0x1A,0x5A,0x9E),
          UP=RGBColor(0x1E,0x7A,0x7A), NOISE=RGBColor(0xB0,0x6A,0x1B), DA=RGBColor(0xC0,0x39,0x2B),
          LOSS=RGBColor(0x6A,0x3D,0x9A), GREY=RGBColor(0x44,0x44,0x44))

def build(bg_path, out_path):
    prs = Presentation()
    prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # SD background (full bleed)
    slide.shapes.add_picture(bg_path, 0, 0, width=Inches(13.333), height=Inches(7.5))

    def add_box(x, y, w, h, kind, text="", fs=10, bold=False, tc=None, fill=True, lw=1.2, alpha=0.85):
        shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        shp.adjustments[0] = 0.08
        if fill:
            shp.fill.solid(); shp.fill.fore_color.rgb = FC[kind]
            # transparency for overlay on bg
            shp.fill.fore_color.brightness = 0  # keep opaque for readability
        else:
            shp.fill.background()
        shp.line.color.rgb = EC[kind]; shp.line.width = Pt(lw)
        shp.shadow.inherit = False
        tf = shp.text_frame; tf.word_wrap = True
        tf.margin_left = Pt(2); tf.margin_right = Pt(2); tf.margin_top = Pt(1); tf.margin_bottom = Pt(1)
        if text:
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
            r = p.add_run(); r.text = text
            r.font.size = Pt(fs); r.font.bold = bold
            r.font.color.rgb = tc if tc else EC[kind]; r.font.name = "Helvetica"
        return shp

    def add_arrow(x1, y1, x2, y2, color=RGBColor(0x44,0x44,0x44), lw=1.4):
        conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
        conn.line.color.rgb = color; conn.line.width = Pt(lw)
        ln = conn.line._get_or_add_ln()
        ln.append(ln.makeelement(qn('a:tailEnd'), {'type':'arrow','w':'med','len':'med'}))
        return conn

    add_box(0.4, 0.2, 12.5, 0.55, "DATA", "SDA-Diff: generative data assimilation for dense storm-surge forecasting", fs=15, bold=True, tc=EC["DATA"])
    # input
    add_box(0.4, 1.0, 2.9, 5.9, "DATA", "", fill=True)
    add_box(0.55, 1.15, 2.6, 0.4, "DATA", "Input", fs=11, bold=True, tc=EC["DATA"])
    add_box(0.55, 1.7, 2.6, 0.75, "DATA", "sparse in-situ gauges\n+ valid mask", fs=8)
    add_box(0.55, 2.55, 2.6, 0.75, "DATA", "ERA5  msl·u10·v10", fs=8)
    add_box(0.55, 3.4, 2.6, 0.75, "DATA", "GTSM surge", fs=8)
    add_box(0.55, 4.3, 2.6, 0.75, "DATA", "c=[B,T,6,H,W]  +  lead L", fs=8, bold=True)
    add_arrow(0.55, 5.2, 3.35, 5.6)
    # model
    add_box(3.5, 1.0, 6.3, 5.9, "CONV", "", fill=True, lw=1.6)
    add_box(3.65, 1.15, 6.0, 0.4, "CONV", "Model — EDM generation + SDA assimilation", fs=11, bold=True, tc=EC["CONV"])
    # training
    add_box(3.7, 1.75, 5.95, 2.1, "CONV", "", fill=False, lw=1.0)
    add_box(3.85, 1.85, 5.7, 0.35, "CONV", "Training — learn to recover the surge field from noise", fs=9, bold=True, tc=EC["CONV"], fill=False)
    add_box(3.9, 2.3, 1.2, 0.7, "DATA", "x₀ target", fs=8)
    add_box(5.2, 2.3, 1.7, 0.7, "NOISE", "xₜ = x₀ + σε\n(noise / drop)", fs=7.5)
    add_arrow(5.1, 2.65, 5.2, 2.65, color=EC["NOISE"])
    add_box(7.0, 2.3, 2.4, 0.7, "CONV", "D_θ(xₜ,σ,c,L)  U-Net\nconv·down·up·attn", fs=7.5, bold=True)
    add_arrow(6.9, 2.65, 7.0, 2.65)
    # inference
    add_box(3.7, 4.1, 5.95, 2.65, "DA", "", fill=False, lw=1.4)
    add_box(3.85, 4.2, 5.7, 0.35, "DA", "Inference — correct the sampling path with new observations", fs=9, bold=True, tc=EC["DA"], fill=False)
    add_box(3.9, 4.65, 1.1, 0.7, "DATA", "noise xₜ", fs=8)
    add_box(5.1, 4.65, 1.3, 0.7, "CONV", "denoise\n(EDM Heun)", fs=7.5)
    add_arrow(5.0, 5.0, 5.1, 5.0)
    add_box(6.5, 4.65, 1.3, 0.7, "CONV", "Tweedie x̂₀", fs=7.5)
    add_arrow(6.4, 5.0, 6.5, 5.0)
    add_box(7.9, 4.65, 1.7, 0.7, "DATA", "obs. y, A, R\n(gauge pixels)", fs=7.5)
    add_box(5.0, 5.55, 4.5, 1.0, "DA", "SDA:  ∇log p(xₜ|y,c) = s_θ + ∇log N(y|Ax̂₀,R)", fs=8.5, bold=True, tc=EC["DA"], lw=1.8)
    add_arrow(7.1, 4.65, 7.25, 5.55, color=EC["DA"]); add_arrow(8.7, 5.0, 7.25, 5.55, color=EC["DA"])
    # output
    add_box(10.0, 1.0, 3.0, 5.9, "DATA", "", fill=True)
    add_box(10.15, 1.15, 2.7, 0.4, "DATA", "Output (ensemble)", fs=11, bold=True, tc=EC["DATA"])
    add_box(10.2, 2.2, 2.6, 0.7, "CONV", "posterior {x₀⁽ⁱ⁾} (N samples)", fs=8)
    add_box(10.2, 3.3, 2.6, 1.0, "UP", "mean · quantiles\nP(surge > h)", fs=8.5)
    add_box(10.2, 4.7, 2.6, 0.8, "UP", "dense surge forecast\n(uncertainty-aware)", fs=8, bold=True)
    add_arrow(9.8, 2.6, 10.0, 2.6); add_arrow(7.25, 5.9, 10.0, 3.6, color=EC["DA"])
    # loss
    add_box(3.7, 7.0, 6.0, 0.45, "LOSS", "L = E[λ(σ) ‖ D_θ(x₀+σε, σ, c, L) − x₀ ‖²]  (masked)", fs=9, bold=True, tc=EC["LOSS"])
    add_arrow(6.5, 3.85, 6.5, 7.0, color=EC["LOSS"])
    prs.save(out_path)
    print("saved", out_path)

for seed in [0, 1, 2]:
    build(f"temp/figures/sd_fig_seed{seed}.png", f"temp/figures/fig1_framework_sd{seed}.pptx")
